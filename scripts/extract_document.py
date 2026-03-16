#!/usr/bin/env python3
"""
KXKM_Clown — Document Text Extraction

Extracts text from various document formats:
- Word (.docx, .doc)
- Excel (.xlsx, .xls)
- PowerPoint (.pptx, .ppt)
- LibreOffice (.odt, .ods, .odp)
- RTF (.rtf)
- EPUB (.epub)

Usage:
  python scripts/extract_document.py --input /path/to/file.docx

Output: JSON on stdout with extracted text.

Install dependencies:
  pip install python-docx openpyxl python-pptx odfpy striprtf EbookLib
"""

import argparse
import json
import os
import sys
import time


def parse_args():
    p = argparse.ArgumentParser(description="KXKM Document Extraction")
    p.add_argument("--input", required=True, help="Path to document file")
    p.add_argument("--max-chars", type=int, default=12000, help="Max chars to extract")
    return p.parse_args()


def extract_docx(path):
    """Extract text from .docx (Word)."""
    from docx import Document
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def extract_xlsx(path):
    """Extract text from .xlsx (Excel)."""
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"[Feuille: {sheet.title}]")
        for row in sheet.iter_rows(max_row=200, values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = "\t".join(cells).strip()
            if line:
                lines.append(line)
    wb.close()
    return "\n".join(lines)


def extract_pptx(path):
    """Extract text from .pptx (PowerPoint)."""
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def extract_odt(path):
    """Extract text from .odt (LibreOffice Writer)."""
    from odf.opendocument import load
    from odf.text import P
    doc = load(path)
    paragraphs = []
    for p in doc.getElementsByType(P):
        text = ""
        for node in p.childNodes:
            if hasattr(node, "data"):
                text += node.data
            elif hasattr(node, "__str__"):
                text += str(node)
        if text.strip():
            paragraphs.append(text.strip())
    return "\n".join(paragraphs)


def extract_rtf(path):
    """Extract text from .rtf."""
    from striprtf.striprtf import rtf_to_text
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return rtf_to_text(f.read())


def extract_epub(path):
    """Extract text from .epub."""
    import ebooklib
    from ebooklib import epub
    from html.parser import HTMLParser

    class TextExtractor(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text = []
        def handle_data(self, data):
            self.text.append(data)

    book = epub.read_epub(path)
    texts = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        parser = TextExtractor()
        parser.feed(item.get_content().decode("utf-8", errors="ignore"))
        texts.extend(parser.text)
    return "\n".join(t.strip() for t in texts if t.strip())


# Map extensions to extractors
EXTRACTORS = {
    ".docx": extract_docx,
    ".doc": extract_docx,  # python-docx handles .doc partially
    ".xlsx": extract_xlsx,
    ".xls": extract_xlsx,  # openpyxl handles .xls partially
    ".pptx": extract_pptx,
    ".ppt": extract_pptx,
    ".odt": extract_odt,
    ".ods": extract_odt,  # similar XML structure
    ".odp": extract_odt,
    ".rtf": extract_rtf,
    ".epub": extract_epub,
}


def main():
    args = parse_args()
    start = time.time()
    result = {"status": "failed", "text": "", "error": None}

    try:
        ext = os.path.splitext(args.input)[1].lower()
        extractor = EXTRACTORS.get(ext)

        if not extractor:
            raise ValueError(f"Format non supporté: {ext}")

        text = extractor(args.input)
        text = text[:args.max_chars]
        duration = time.time() - start

        result = {
            "status": "completed",
            "text": text,
            "format": ext,
            "chars": len(text),
            "duration": round(duration, 2),
        }
        print(f"[extract] {ext} → {len(text)} chars in {duration:.1f}s", file=sys.stderr)

    except ImportError as e:
        module = str(e).split("'")[-2] if "'" in str(e) else str(e)
        result["error"] = f"Module manquant: {module}. Install: pip install python-docx openpyxl python-pptx odfpy striprtf EbookLib"
        print(f"[extract] ERROR: {result['error']}", file=sys.stderr)
    except Exception as e:
        result["error"] = str(e)
        print(f"[extract] ERROR: {e}", file=sys.stderr)

    print(json.dumps(result))


if __name__ == "__main__":
    main()
